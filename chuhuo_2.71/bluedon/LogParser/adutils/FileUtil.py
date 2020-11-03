# -*-coding:utf-8 -*-

import os
import xlrd  # for xls file
import re
import traceback
import logging
import importlib
import mimetypes
import magic
import shutil
import subprocess
from pptx import Presentation

try:
	_magic = magic.open(magic.MIME_TYPE)
	_magic.load()
except:
	_magic = None


# 取出html格式
def clean_html(html):
	"""
	Copied from NLTK package.
	Remove HTML markup from the given string.

	:param html: the HTML string to be cleaned
	:type html: str
	:rtype: str
	"""

	# First we remove inline JavaScript/CSS:
	cleaned = re.sub(r"(?is)<(script|style).*?>.*?(</\1>)", "", html.strip())
	# Then we remove html comments. This has to be done before removing regular
	# tags since comments can contain '>' characters.
	cleaned = re.sub(r"(?s)<!--(.*?)-->[\n]?", "", cleaned)
	# Next we can remove the remaining tags:
	cleaned = re.sub(r"(?s)<.*?>", " ", cleaned)
	cleaned = re.sub(r"\n\n", " ", cleaned)
	cleaned = re.sub(r"\n\n", " ", cleaned)
	cleaned = re.sub(r"\n\n", " ", cleaned)
	# Finally, we deal with whitespace
	cleaned = re.sub(r"&nbsp;", " ", cleaned)
	cleaned = re.sub(r"  ", " ", cleaned)
	cleaned = re.sub(r"  ", " ", cleaned)
	return cleaned.strip()


def _check_mime(filename):
	""" This function checks the magic of a file referenced by ``filename``
	returns possible filename extension(s) based on detected mimetype
	(could be a string or a list)
	"""
	ext = []
	try:
		if hasattr(magic, 'from_file'):
			mimetype = magic.from_file(filename, mime=True)
		elif _magic:
			mimetype = _magic.file(filename)

		if mimetype:
			mimetype = mimetype.split(";")[0]
		if not mimetype:
			mimetype = None
	except:
		mimetype = None
	ext = mimetypes.guess_all_extensions(mimetype)
	return ext


def textract_process(filename, encoding='utf_8', mime=None, **kwargs):
	"""This is the core function used for extracting text. It routes the
	``filename`` to the appropriate parser and returns the extracted
	text as a byte-string encoded with ``encoding``.
	"""
	if not os.path.exists(filename):
		logging.debug('%s not exist')
		return
	_, ext1 = os.path.splitext(filename)
	ext1 = ext1.lower()

	ext3 = []
	ext2 = _check_mime(filename)

	if isinstance(ext2, list):
		for i in ext2:
			ext3.append(i)
	else:
		ext3.append(ext2)

	if isinstance(ext1, list):
		for i in ext1:
			ext3.append(i)
	else:
		ext3.append(ext1)
	for ext in ext3:
		logging.debug("Processing ext %s" % ext)
		if ext:
			rel_module = ext + '_parser'
			module_name = rel_module[1:]
			try:
				filetype_module = importlib.import_module(rel_module,
														  'textract.parsers')
			except Exception as e:
				continue
			parser = filetype_module.Parser()
			ret = parser.process(filename, encoding, **kwargs)
			if ret:
				return ret
			else:
				continue
	logging.debug('textract %s fail', os.path.basename(filename))


def _textract_cmd(docfile, txtfile):
	ret = -1
	try:
		logging.debug("texract %s > %s",
					  os.path.basename(docfile),
					  os.path.basename(txtfile))
		cmd = 'textract "%s" > "%s"' % (docfile, txtfile)
		# ret = os.system(cmd)
		p = subprocess.Popen(md,
							 shell=True,
							 stdout=subprocess.PIPE,
							 stderr=subprocess.STDOUT,
							 close_fds=True)
		p.communicate()
	except:
		logging.debug("try textract commond fail")

	if ret != 0:
		try:
			logging.debug("catdoc %s > %s",
						  os.path.basename(docfile),
						  os.path.basename(txtfile))
			ret = os.system('catdoc "%s" > "%s"' % (docfile, txtfile))
		except:
			logging.debug("try catdoc fail")
	return ret


def _textract_file(in_file, out_file):
	try:
		txt = textract_process(in_file)
		if txt:
			logging.debug('textract_process %s to %s successfully',
						  os.path.basename(in_file),
						  os.path.basename(out_file))
			with open(out_file, 'w') as f:
				f.write(txt)
			return 0
	except Exception as e:
		logging.warning('textract_process fail')
	return _textract_cmd(in_file, out_file)


def doc2txt(doc_file, txt_file):
	"""将doc文本转化为txt 成功返回0，失败返回1

	使用magic模块检测文件类型，防止doc与zip文件混淆，textract程序无法
	解析将word文件后缀名改为zip的文件，先尝试用代码的方式转换word为text，
	如果失败在尝试用命令行方式
	"""
	return _textract_file(doc_file, txt_file)


def _convert_by_pptx(ppt_file):
	"""pptx 转文本，只支持pptx 2007之后版本"""
	try:
		prs = Presentation(ppt_file)
		# text_runs will be populated with a list of strings,
		# one for each text run in presentation
		text_runs = []

		for slide in prs.slides:
			for shape in slide.shapes:
				if not shape.has_text_frame:
					continue
				for paragraph in shape.text_frame.paragraphs:
					for run in paragraph.runs:
						text_runs.append(run.text)
		return u'\n'.join(text_runs).encode('utf8')
	except Exception, e:
		return None


def ppt2txt(ppt_file, txt_file):
	"""将ppt转换为txt"""
	text = _convert_by_pptx(ppt_file)
	if text == None:
		logging.debug("convert %s by pptx fail, try to textract",
					  os.path.basename(ppt_file))
		# 调用 textract 尝试进行转换
		return _textract_file(ppt_file, txt_file)
	else:
		with open(txt_file, 'w') as f:
			f.write(text)
		return 0


def xls2txt(xls_file, txt_file):
	data = xlrd.open_workbook(xls_file)
	fileobject = open(txt_file, 'w')
	sheets = data.sheets()
	sheet_names = data.sheet_names()
	length = len(sheets)
	list = []
	for i in range(len(sheets)):
		rows = sheets[i].nrows
		list = []
		for row in range(0, rows):
			value = sheets[i].row_values(row)
			list.append(value)

		if list:
			fileobject.write(sheet_names[i])
			fileobject.write('\n')
			for i in range(len(list)):
				for j in range(len(list[i])):
					fileobject.write(str(list[i][j]))
					fileobject.write("\t")
				fileobject.write('\n')
			fileobject.write('\n')
	fileobject.close()

	if os.path.getsize(txt_file) > 0:
		return 0
	os.remove(txt_file)
	return 1


def rmdir(src):
	'''delete files and folders'''
	shutil.rmtree(src)


if __name__ == '__main__':
	import sys

	ppt2txt(sys.argv[1], sys.argv[2])